"""
FINQUANT-NEXUS v4 — PPT Generator
Run: python3 generate_ppt.py
Output: FINQUANT_NEXUS_Review2.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Color Palette ────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)   # dark navy
BG_CARD    = RGBColor(0x1E, 0x29, 0x3B)   # slate card
BG_CARD2   = RGBColor(0x0F, 0x28, 0x1E)   # dark green tint card
ACCENT_G   = RGBColor(0x10, 0xB9, 0x81)   # emerald green
ACCENT_B   = RGBColor(0x38, 0xBD, 0xF8)   # sky blue
ACCENT_A   = RGBColor(0xFB, 0xBF, 0x24)   # amber
ACCENT_R   = RGBColor(0xF8, 0x71, 0x71)   # soft red
WHITE      = RGBColor(0xF8, 0xFA, 0xFC)   # near white
GRAY       = RGBColor(0x94, 0xA3, 0xB8)   # slate gray
DIM        = RGBColor(0x64, 0x74, 0x8B)   # dimmer gray

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height

# ─── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, color=None):
    color = color or BG_DARK
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.width = 0
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    color = color or WHITE
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_para(tf, text, size=16, bold=False, color=None,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    color = color or WHITE
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def accent_bar(slide, color=None, height=Inches(0.07)):
    color = color or ACCENT_G
    add_rect(slide, 0, 0, W, height, color)

def slide_title(slide, title, subtitle=None, title_color=None):
    accent_bar(slide)
    title_color = title_color or ACCENT_G
    add_text(slide, title,
             Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.65),
             size=28, bold=True, color=title_color)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.75), Inches(12), Inches(0.4),
                 size=15, color=GRAY)

def divider(slide, y, color=None):
    color = color or DIM
    add_rect(slide, Inches(0.5), y, Inches(12.33), Inches(0.01), color)


# ─── SLIDE BUILDERS ───────────────────────────────────────────────────────────

def slide_01_title(prs):
    s = blank_slide(prs)
    fill_bg(s)

    # top accent bar (thicker)
    add_rect(s, 0, 0, W, Inches(0.12), ACCENT_G)

    # left accent stripe
    add_rect(s, 0, 0, Inches(0.08), H, ACCENT_G)

    # main title
    add_text(s, "FINQUANT-NEXUS", Inches(1.2), Inches(1.6), Inches(11), Inches(1.3),
             size=54, bold=True, color=ACCENT_G, align=PP_ALIGN.CENTER)

    add_text(s, "v4",
             Inches(9.6), Inches(1.65), Inches(1.5), Inches(0.6),
             size=22, bold=False, color=ACCENT_A)

    # tagline
    add_text(s,
             "Self-Optimizing Federated Portfolio Intelligence Platform for NIFTY 50",
             Inches(1.2), Inches(2.85), Inches(11), Inches(0.6),
             size=20, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    divider(s, Inches(3.6))

    # meta info row
    add_text(s, "Praveen Rawal",
             Inches(1.2), Inches(3.75), Inches(5), Inches(0.45),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "M.Tech – Data Science & Machine Learning  |  2nd Year",
             Inches(1.2), Inches(4.2), Inches(11), Inches(0.4),
             size=16, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, "Project Review – 2",
             Inches(1.2), Inches(4.65), Inches(11), Inches(0.4),
             size=15, color=ACCENT_A, align=PP_ALIGN.CENTER)

    # bottom badge row (tech pills)
    pills = ["Reinforcement Learning", "Graph Neural Networks",
             "Federated Learning", "FinBERT Sentiment", "NIFTY 50"]
    px = Inches(0.6)
    for pill in pills:
        pw = Inches(2.3)
        add_rect(s, px, Inches(6.6), pw, Inches(0.45), BG_CARD)
        add_text(s, pill, px + Inches(0.1), Inches(6.62), pw - Inches(0.2), Inches(0.42),
                 size=11, color=ACCENT_G, align=PP_ALIGN.CENTER)
        px += Inches(2.45)

    # bottom bar
    add_rect(s, 0, H - Inches(0.08), W, Inches(0.08), ACCENT_G)
    return s


def slide_02_agenda(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_title(s, "Agenda", "What we'll cover today")

    items = [
        ("01", "Problem Statement",     "Why existing portfolio systems fall short"),
        ("02", "Solution Overview",     "Architecture & approach of FINQUANT-NEXUS"),
        ("03", "Dataset",               "NIFTY 50 data, sources, and preprocessing"),
        ("04", "Tabs & Techniques",     "7 interactive modules — deep dive"),
        ("05", "Tech Stack",            "Frontend, backend, and ML libraries"),
        ("06", "Key Results",           "Performance metrics achieved"),
        ("07", "Future Work",           "Remaining 20% — roadmap"),
        ("08", "Conclusion",            "Summary and impact"),
    ]

    col1_x, col2_x = Inches(0.5), Inches(6.8)
    row_h = Inches(0.72)
    start_y = Inches(1.25)

    for i, (num, title, desc) in enumerate(items):
        col_x = col1_x if i < 4 else col2_x
        row_y = start_y + (i % 4) * row_h

        add_rect(s, col_x, row_y, Inches(0.55), Inches(0.52), ACCENT_G)
        add_text(s, num, col_x + Inches(0.04), row_y + Inches(0.05),
                 Inches(0.47), Inches(0.42), size=14, bold=True,
                 color=BG_DARK, align=PP_ALIGN.CENTER)
        add_text(s, title,
                 col_x + Inches(0.65), row_y, Inches(5.5), Inches(0.3),
                 size=15, bold=True, color=WHITE)
        add_text(s, desc,
                 col_x + Inches(0.65), row_y + Inches(0.28), Inches(5.5), Inches(0.28),
                 size=11, color=GRAY)
    return s


def slide_03_problem(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_title(s, "Problem Statement", "Gaps in modern portfolio management")

    problems = [
        (ACCENT_R, "Static Optimization",
         "Markowitz & MPT assume fixed correlations & normal returns. Markets are dynamic, non-stationary, and regime-changing."),
        (ACCENT_A, "Ignored News & Sentiment",
         "Price-only models miss news signals. Earnings surprises, macroeconomic announcements move stocks within hours — before price reflects it."),
        (ACCENT_B, "No Privacy in Collaboration",
         "Fund houses cannot share raw trade data (SEBI / GDPR). Collaborative model training was impossible — until federated learning."),
        (ACCENT_G, "Missing Stock Relationships",
         "Stocks treated as independent entities. Supply-chain links, sector co-movements, and correlations are completely ignored."),
        (RGBColor(0xA7, 0x8B, 0xFA), "No Stress Awareness",
         "Most models optimize for returns only. No testing against real crash scenarios: 2008 crisis, COVID crash, flash crashes."),
    ]

    for i, (color, title, desc) in enumerate(problems):
        row = i // 2 if i < 4 else None
        if i < 4:
            col = i % 2
            bx = Inches(0.4 + col * 6.4)
            by = Inches(1.3 + row * 2.75)
            bw = Inches(6.1)
            bh = Inches(2.45)
        else:
            bx = Inches(3.6)
            by = Inches(6.82)
            bw = Inches(6.1)
            bh = Inches(0.55)

        if i < 4:
            add_rect(s, bx, by, bw, bh, BG_CARD)
            add_rect(s, bx, by, Inches(0.08), bh, color)
            add_text(s, title,
                     bx + Inches(0.2), by + Inches(0.15),
                     bw - Inches(0.3), Inches(0.38),
                     size=16, bold=True, color=color)
            add_text(s, desc,
                     bx + Inches(0.2), by + Inches(0.55),
                     bw - Inches(0.4), Inches(1.7),
                     size=12, color=GRAY)
        else:
            add_rect(s, bx, Inches(1.3 + 2 * 2.75), bw, Inches(2.45), BG_CARD)
            add_rect(s, bx, Inches(1.3 + 2 * 2.75), Inches(0.08), Inches(2.45), color)
            add_text(s, title,
                     bx + Inches(0.2), Inches(1.3 + 2 * 2.75) + Inches(0.15),
                     bw - Inches(0.3), Inches(0.38),
                     size=16, bold=True, color=color)
            add_text(s, desc,
                     bx + Inches(0.2), Inches(1.3 + 2 * 2.75) + Inches(0.55),
                     bw - Inches(0.4), Inches(1.7),
                     size=12, color=GRAY)
    return s


def slide_04_solution(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_title(s, "Solution Overview", "FINQUANT-NEXUS v4 — 3-layer intelligent pipeline")

    layers = [
        (ACCENT_B,  "Layer 1 — Data Intelligence",
         ["Yahoo Finance: NIFTY 50 prices (2015–2025, ~10 years)",
          "Google News RSS: live headlines for 44 stocks",
          "21 technical indicators: RSI, MACD, Bollinger Bands, ATR…"]),
        (ACCENT_A,  "Layer 2 — Model Intelligence",
         ["T-GAT: Temporal Graph Attention Network — 3 edge types",
          "5 RL Agents (PPO · SAC · TD3 · A2C · DDPG) → Ensemble",
          "FinBERT: domain-specific financial NLP (ProsusAI/finbert)"]),
        (ACCENT_G,  "Layer 3 — Federated Intelligence",
         ["4 sector clients: Banking · Finance · IT · Others",
          "FedProx aggregation — handles non-IID sector data",
          "DP-SGD privacy: ε=8.0, δ=10⁻⁵ — gradient noise injection"]),
    ]

    for i, (color, title, points) in enumerate(layers):
        bx = Inches(0.35 + i * 4.33)
        by = Inches(1.3)
        bw = Inches(4.1)
        bh = Inches(5.6)
        add_rect(s, bx, by, bw, bh, BG_CARD)
        add_rect(s, bx, by, bw, Inches(0.06), color)

        # Layer number circle
        add_rect(s, bx + Inches(1.55), by + Inches(0.15), Inches(1.0), Inches(0.55), color)
        add_text(s, f"0{i+1}", bx + Inches(1.55), by + Inches(0.15),
                 Inches(1.0), Inches(0.55), size=18, bold=True,
                 color=BG_DARK, align=PP_ALIGN.CENTER)

        add_text(s, title,
                 bx + Inches(0.15), by + Inches(0.85),
                 bw - Inches(0.3), Inches(0.5),
                 size=14, bold=True, color=color)

        for j, pt in enumerate(points):
            add_text(s, f"→  {pt}",
                     bx + Inches(0.15), by + Inches(1.5 + j * 1.25),
                     bw - Inches(0.25), Inches(1.1),
                     size=12, color=WHITE)

    # bottom label
    add_text(s, "Final Output: Privacy-preserving, ensemble-optimized NIFTY 50 portfolio allocation",
             Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35),
             size=13, italic=True, color=ACCENT_G, align=PP_ALIGN.CENTER)
    return s


def slide_05_dataset(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_title(s, "Dataset & Data Sources", "Real-world financial data — no synthetic benchmarks")

    # Left block
    add_rect(s, Inches(0.4), Inches(1.3), Inches(5.9), Inches(5.6), BG_CARD)
    add_text(s, "Price Data — Yahoo Finance (yfinance)",
             Inches(0.6), Inches(1.4), Inches(5.5), Inches(0.45),
             size=16, bold=True, color=ACCENT_B)

    rows = [
        ("Universe",    "NIFTY 50 — India's top 50 large-cap stocks"),
        ("Coverage",    "44 stocks with complete data (of 50)"),
        ("Period",      "January 2015 – December 2025 (~10 years)"),
        ("Granularity", "Daily OHLCV + adjusted close"),
        ("Features",    "21 technical indicators per stock per day"),
        ("Split",       "Train: 2015–2022  |  Validate: 2022–2025"),
        ("Cache",       "all_close_prices.csv — avoids repeated API calls"),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(1.95 + i * 0.65)
        add_text(s, k, Inches(0.6), y, Inches(1.6), Inches(0.5),
                 size=12, bold=True, color=ACCENT_B)
        add_text(s, v, Inches(2.3), y, Inches(3.8), Inches(0.5),
                 size=12, color=WHITE)

    # Right block
    add_rect(s, Inches(6.7), Inches(1.3), Inches(6.2), Inches(2.55), BG_CARD)
    add_text(s, "News Data — Google News RSS",
             Inches(6.9), Inches(1.4), Inches(5.8), Inches(0.45),
             size=16, bold=True, color=ACCENT_A)
    news_rows = [
        ("Queries",   "22 (20 stock-specific + 2 market-wide)"),
        ("Fetch",     "Live HTTP requests — real-time headlines"),
        ("NLP Model", "FinBERT (ProsusAI) — financial domain BERT"),
        ("Output",    "Score [-1,+1] + positive/neutral/negative %"),
        ("Refresh",   "Auto-refresh every 3 minutes in dashboard"),
    ]
    for i, (k, v) in enumerate(news_rows):
        y = Inches(1.95 + i * 0.68)
        add_text(s, k, Inches(6.9), y, Inches(1.5), Inches(0.5),
                 size=12, bold=True, color=ACCENT_A)
        add_text(s, v, Inches(8.5), y, Inches(4.2), Inches(0.5),
                 size=12, color=WHITE)

    # Right bottom — 21 indicators
    add_rect(s, Inches(6.7), Inches(4.05), Inches(6.2), Inches(2.8), BG_CARD)
    add_text(s, "21 Technical Indicators (Feature Engineering)",
             Inches(6.9), Inches(4.15), Inches(5.8), Inches(0.45),
             size=15, bold=True, color=ACCENT_G)
    indicators = "RSI · MACD · MACD Signal · Bollinger Upper/Mid/Lower · ATR · OBV · CCI · Williams %R · Stochastic K/D · EMA-5/10/20/50 · SMA-20 · ROC · MFI · VWAP"
    add_text(s, indicators,
             Inches(6.9), Inches(4.7), Inches(5.8), Inches(1.9),
             size=11, color=GRAY)
    return s


def tab_slide(prs, num, title, route, accent, what_it_does,
              features, technique_title, techniques, impressive):
    s = blank_slide(prs)
    fill_bg(s)

    # Colored top bar
    add_rect(s, 0, 0, W, Inches(0.1), accent)

    # Tab badge
    add_rect(s, Inches(0.4), Inches(0.15), Inches(0.6), Inches(0.5), accent)
    add_text(s, str(num), Inches(0.4), Inches(0.15), Inches(0.6), Inches(0.5),
             size=16, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

    # Title & route
    add_text(s, title,
             Inches(1.1), Inches(0.15), Inches(8), Inches(0.5),
             size=24, bold=True, color=accent)
    add_text(s, f"Route: {route}",
             Inches(1.1), Inches(0.65), Inches(4), Inches(0.35),
             size=12, color=DIM)

    divider(s, Inches(0.95))

    # LEFT COLUMN — What it does + Features
    add_rect(s, Inches(0.4), Inches(1.1), Inches(5.7), Inches(5.9), BG_CARD)
    add_text(s, "What It Does",
             Inches(0.6), Inches(1.2), Inches(5.3), Inches(0.35),
             size=14, bold=True, color=accent)
    add_text(s, what_it_does,
             Inches(0.6), Inches(1.6), Inches(5.3), Inches(1.0),
             size=12, color=WHITE)

    add_text(s, "Key Features",
             Inches(0.6), Inches(2.65), Inches(5.3), Inches(0.35),
             size=14, bold=True, color=accent)
    for i, feat in enumerate(features):
        add_text(s, f"✦  {feat}",
                 Inches(0.65), Inches(3.1 + i * 0.56),
                 Inches(5.2), Inches(0.5),
                 size=11.5, color=WHITE)

    # RIGHT COLUMN — Techniques
    add_rect(s, Inches(6.5), Inches(1.1), Inches(6.4), Inches(5.9), BG_CARD)
    add_text(s, technique_title,
             Inches(6.7), Inches(1.2), Inches(6.0), Inches(0.35),
             size=14, bold=True, color=ACCENT_A)
    for i, tech in enumerate(techniques):
        add_text(s, f"◆  {tech}",
                 Inches(6.7), Inches(1.65 + i * 0.62),
                 Inches(6.0), Inches(0.56),
                 size=11.5, color=WHITE)

    # BOTTOM bar — Why impressive
    add_rect(s, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3), BG_CARD2)
    add_text(s, f"★  {impressive}",
             Inches(0.6), Inches(7.1), Inches(12.1), Inches(0.3),
             size=11, color=ACCENT_G)
    return s


def slide_06_to_12_tabs(prs):
    tabs = [
        (1, "Portfolio Overview", "/",
         ACCENT_G,
         "Final optimized output of the system — shows performance metrics, all 44 NIFTY holdings, sector allocation, and portfolio vs NIFTY benchmark chart.",
         ["6 metric cards: Sharpe, Sortino, Return, Volatility, Drawdown (with badge thresholds)",
          "Holdings table: 44 stocks, weight%, return%, color-coded sector dots",
          "Sector allocation: horizontal BarChart (11 sectors)",
          "Stock click → 60-day price chart, 52W range, risk grid",
          "Investment Simulator: ₹ input → per-stock P&L + Growth Chart",
          "3-way benchmark: RL Portfolio vs NIFTY 50 vs FD (7%)"],
         "APIs & Techniques",
         ["/api/portfolio-summary  →  Sharpe, holdings, sector weights",
          "/api/stock/{ticker}    →  60-day history, 52W H/L, daily Δ",
          "/api/portfolio-growth  →  Time-series: Portfolio vs NIFTY vs FD",
          "/api/refresh-data      →  Live yfinance pull (Yahoo Finance)",
          "Sharpe = (Return − 7%) / σ  |  Sortino uses downside σ only",
          "Investment Sim: softmax weights × invested amount × 1Y return",
          "Recharts: BarChart, AreaChart, LineChart (3-series growth)"],
         "Real 10-year NIFTY data + interactive simulator with 3-way benchmark comparison"),

        (2, "RL Agent", "/rl",
         ACCENT_B,
         "Trains and compares 5 reinforcement learning algorithms on NIFTY 50 portfolio allocation. Ensemble combines all 5 for best risk-adjusted returns.",
         ["6 algo selector: PPO · SAC · TD3 · A2C · DDPG · Ensemble",
          "Comparison table: Sharpe, Sortino, Return, Volatility, Drawdown",
          "3-tab charts: Training Rewards / Cumulative Returns / Weights",
          "Sector allocation donut per algorithm",
          "Return contribution list: top 15 stocks with animated bars",
          "Constraints panel: position limits, stop-loss, circuit breaker"],
         "RL Algorithms & Design",
         ["Stable-Baselines3: PPO (clipped PG), SAC (entropy reg.), TD3 (twin critics)",
          "Environment: 21 features × 44 stocks state space",
          "Reward: Sharpe ratio at each rebalancing step",
          "Constraints: max pos 20%, stop-loss -5%, tx cost 0.1%",
          "Train: 2019–2022 | Validate: 2022–2023 (unseen split)",
          "Ensemble = weighted average of 5 agent outputs",
          "Recharts: AreaChart (rewards), LineChart (returns), PieChart (sector)"],
         "5 independent algorithms + ensemble — diversity reduces variance, outperforms any single agent"),

        (3, "Stress Testing", "/stress",
         ACCENT_R,
         "Tests portfolio resilience under 4 historical crash scenarios using Monte Carlo simulation. Computes VaR and CVaR — industry-standard risk metrics.",
         ["4 crash scenarios: Normal / 2008 Crisis / COVID Crash / Flash Crash",
          "Metrics: VaR 95%, CVaR 95%, Survival Rate, Mean Return per scenario",
          "Monte Carlo fan chart: 30 simulated paths over 60 days",
          "Color-coded risk severity (green → dark red)",
          "User-configurable: n_stocks (2–47), n_simulations (100–50K)",
          "Expandable metric cards with formula explanations"],
         "Risk Computation Methods",
         ["/api/stress-test  POST {n_stocks, n_simulations}",
          "Monte Carlo: multivariate normal returns × correlation matrix",
          "2008 params: 3.5% daily vol + 30% correlation spike",
          "COVID params: 5.0% daily vol + 40% correlation spike",
          "Flash params: 8.0% daily vol, 5-day extreme event",
          "VaR 95% = 5th percentile of all simulation final values",
          "CVaR 95% = E[loss | loss > VaR] — coherent risk measure (Basel III)"],
         "VaR + CVaR with 4 real crisis calibrations — portfolio survives 82%+ of COVID simulations"),

        (4, "Federated Learning", "/fl",
         RGBColor(0xA7, 0x8B, 0xFA),
         "Privacy-preserving collaborative learning across 4 sector clients. Each client trains locally and shares only gradients — never raw portfolio data.",
         ["4 sector clients: Banking(15) · Finance(10) · IT(6) · Others(19)",
          "Strategy comparison: FedProx vs FedAvg convergence (50 rounds)",
          "Privacy dashboard: ε=8.0, δ=10⁻⁵ (DP-SGD)",
          "Convergence chart: 6 lines (global + 4 clients)",
          "Fairness comparison: Sharpe with FL vs without FL per client",
          "Non-IID sector split — realistic, not artificial IID assumption"],
         "FL Algorithms & Privacy",
         ["FedAvg: global = Σ (n_i/N) × local_weights_i",
          "FedProx: adds proximal term μ||w − w_global||² to local loss",
          "DP-SGD: clip gradients to norm C, add N(0, σ²C²) noise",
          "Privacy calibration: (ε=8.0, δ=10⁻⁵) — mathematical guarantee",
          "Custom implementation in PyTorch (src/federated/)",
          "4 clients are non-IID by design (banking ≠ IT behavior)",
          "All 4 clients show improved Sharpe with FL vs isolated training"],
         "DP-SGD privacy from scratch — solves real regulatory constraint (SEBI/GDPR data sharing)"),

        (5, "Sentiment Analysis", "/sentiment",
         ACCENT_A,
         "Real-time financial news analysis using FinBERT. Live Google News headlines processed every 3 minutes — sentiment adjusts portfolio weights automatically.",
         ["Live badge + auto-refresh every 3 min (no page reload)",
          "22 RSS queries → batch FinBERT → per-headline scores",
          "3-tab view: News Feed / Portfolio Impact / Sector Breakdown",
          "Manual text analysis box — type any text, get instant score",
          "High-impact alerts: |score| > 0.3 filtered automatically",
          "Trend chart persisted in localStorage across page refreshes"],
         "NLP Pipeline",
         ["Model: ProsusAI/finbert — BERT fine-tuned on financial text",
          "Output: score ∈ [-1,+1] + positive/neutral/negative probabilities",
          "Weight adjust: w_adj = w_base × (1 + 0.1 × sentiment_signal)",
          "Concurrent RSS fetching: 22 async HTTP requests (feedparser)",
          "Batch inference: predict_batch(texts, batch_size=16) — GPU/CPU",
          "Cache TTL: 3 min — avoids re-running FinBERT on same headlines",
          "New headlines tracked via Set<string> → '+N New' badge"],
         "Real live headlines + FinBERT scores directly linked to portfolio weight adjustments"),

        (6, "Graph Visualization", "/graph",
         ACCENT_G,
         "Interactive force-directed network of 44 NIFTY stocks. Demonstrates T-GAT's input graph — 3 edge types: sector, supply-chain, and price correlation.",
         ["Force-directed physics: 150 frames, then stabilizes",
          "44 nodes: size=degree, color=sector (10 unique colors)",
          "3 edge types (toggleable): Sector · Supply Chain · Correlation",
          "Click any node → right panel: sector, connections, neighbors",
          "Moving data packets animate along active edges",
          "Graph stats: density, avg degree, edge counts per type"],
         "Graph Construction & T-GAT",
         ["Sector edges: same-NIFTY-sector stocks fully connected (~160)",
          "Supply-chain: 30 manually curated B2B relationships (~54 edges)",
          "Correlation edges: |Pearson(60-day)| > 0.4 (dynamic)",
          "T-GAT: separate GATConv per edge type → attention weights",
          "GRU layer: temporal encoding of embedding sequence (t-k…t)",
          "Output: 64-dim embedding per stock (vs 21-dim raw features)",
          "Custom physics engine in TypeScript — no D3.js dependency"],
         "3 edge types + T-GAT attention — captures dynamics no correlation matrix can represent"),

        (7, "Pipeline Workflow", "/workflow",
         ACCENT_B,
         "Animated end-to-end system visualization. 15 component nodes, 19 connections, 9 sequential stages — no API dependency, always available.",
         ["15 nodes spanning all system layers (Data → RL → FL → Output)",
          "19 Bezier edges with animated data packet flow",
          "9-stage sequential animation (1.8s per stage)",
          "Play / Pause / Reset / Stage-jump controls",
          "Click any node → technical spec (inputs, outputs, details)",
          "System stats panel: 246 tests, ε=8 privacy, 44 stocks"],
         "Animation & Architecture",
         ["100% client-side — no API calls, always works",
          "SVG native rendering: nodes, edges, arrowheads",
          "CSS flowDash keyframe — dashed edge animation (GPU)",
          "SVG animateMotion — white packets travel along paths",
          "Framer Motion: stage label fade + slide (0.3s transitions)",
          "Stage sequence: Data→Features→Graph→T-GAT→RL Env→5 Agents→Ensemble→FL→Output",
          "Progress stepper: 9 dots, green fill on completion"],
         "Interactive system architecture diagram — replaces static slide, works without backend"),
    ]

    for t in tabs:
        tab_slide(prs, *t)


def slide_13_techstack(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_title(s, "Tech Stack", "Full-stack modern ML platform")

    sections = [
        (ACCENT_B, "Frontend", Inches(0.35), [
            "React 19 + TypeScript",
            "Vite 8  (build tool)",
            "Tailwind CSS v4",
            "Recharts 3.8  (all charts)",
            "Framer Motion 12  (animations)",
            "TanStack Query 5  (API state)",
            "Zustand 5  (global state)",
            "React Router v7",
        ]),
        (ACCENT_G, "Backend / API", Inches(3.75), [
            "FastAPI 0.104+  (REST API)",
            "Uvicorn  (ASGI server)",
            "Pydantic v2  (validation)",
            "Python 3.11",
            "17 API endpoints",
            "CSV cache layer",
            "120s timeout (heavy routes)",
            "CORS + error middleware",
        ]),
        (ACCENT_A, "ML / Deep Learning", Inches(7.15), [
            "PyTorch 2.1+",
            "PyTorch Geometric 2.4+  (GNNs)",
            "Stable-Baselines3 2.2+  (RL)",
            "Transformers 4.35+  (FinBERT)",
            "Gymnasium 0.29+  (RL env)",
            "NumPy · Pandas · SciPy · Sklearn",
            "yfinance  (price data)",
            "feedparser  (news RSS)",
        ]),
        (ACCENT_R, "Testing & Infra", Inches(10.55), [
            "pytest 7.4+",
            "httpx  (API testing)",
            "246 test cases  (14 files)",
            "Docker + docker-compose",
            "PostgreSQL 16  (optional)",
            "SQLAlchemy 2.0+  (ORM)",
            "W&B  (experiment tracking)",
            "configs/base.yaml  (hyperparams)",
        ]),
    ]

    for color, title, x, items in sections:
        add_rect(s, x, Inches(1.25), Inches(3.1), Inches(5.9), BG_CARD)
        add_rect(s, x, Inches(1.25), Inches(3.1), Inches(0.06), color)
        add_text(s, title,
                 x + Inches(0.15), Inches(1.35),
                 Inches(2.8), Inches(0.4),
                 size=15, bold=True, color=color)
        for i, item in enumerate(items):
            add_text(s, f"• {item}",
                     x + Inches(0.15), Inches(1.85 + i * 0.6),
                     Inches(2.8), Inches(0.55),
                     size=11.5, color=WHITE)
    return s


def slide_14_results(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_title(s, "Key Results & Metrics", "Validation on real NIFTY 50 data — 2022–2025")

    metrics = [
        (ACCENT_G,  "1.87",   "Sharpe Ratio",      "EXCEPTIONAL\n(threshold >1.5)"),
        (ACCENT_B,  "28.4%",  "Annual Return",      "Validated on\n2024–2025 data"),
        (ACCENT_A,  "14.2%",  "Volatility",         "MODERATE RISK\n(controlled)"),
        (ACCENT_R,  "-11.3%", "Max Drawdown",       "ACCEPTABLE\n(threshold <20%)"),
        (ACCENT_G,  "82%",    "Survival Rate",      "Under COVID\ncrash scenario"),
        (RGBColor(0xA7,0x8B,0xFA), "ε=8.0", "FL Privacy Budget",
         "DP-SGD guarantee\n(δ=10⁻⁵)"),
    ]

    for i, (color, val, label, note) in enumerate(metrics):
        col = i % 3
        row = i // 3
        bx = Inches(0.4 + col * 4.3)
        by = Inches(1.3 + row * 2.8)
        bw = Inches(4.0)
        bh = Inches(2.5)
        add_rect(s, bx, by, bw, bh, BG_CARD)
        add_rect(s, bx, by, bw, Inches(0.06), color)
        add_text(s, val,
                 bx, by + Inches(0.3), bw, Inches(0.9),
                 size=38, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, label,
                 bx, by + Inches(1.25), bw, Inches(0.4),
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, note,
                 bx, by + Inches(1.7), bw, Inches(0.65),
                 size=11, color=GRAY, align=PP_ALIGN.CENTER)

    add_text(s, "246 / 246 Tests Passing  ·  ~20,000 lines of code  ·  44 NIFTY stocks  ·  500 RL episodes",
             Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.35),
             size=12, color=DIM, align=PP_ALIGN.CENTER)
    return s


def slide_15_future(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_title(s, "Future Work — Remaining 20%",
                "Deployment maturity, production readiness, and advanced tooling")

    items = [
        (ACCENT_G,  "HIGH",   "Portfolio Growth What-If Tool",
         "Backend /api/portfolio-growth is ready. Frontend panel pending — time-slider for 1M–5Y projection with Portfolio vs NIFTY vs FD comparison."),
        (ACCENT_B,  "HIGH",   "Settings & Configuration Page",
         "Sidebar button exists but has no route. Will enable: RL algorithm selection, refresh interval control, CSV/PDF export, alert thresholds."),
        (ACCENT_A,  "MEDIUM", "PostgreSQL Persistence Layer",
         "SQLAlchemy ORM schema is defined. Wiring actual DB read/write will enable multi-session portfolio history, multi-user support, and audit trail."),
        (RGBColor(0xA7,0x8B,0xFA), "MEDIUM", "JWT Auth & Role-Based Access",
         "No login system currently. Planned: JWT authentication, read-only vs admin roles, per-user portfolio strategy isolation."),
        (ACCENT_R,  "LOW",    "W&B Experiment Tracker Integration",
         "Backend already supports wandb logging (optional). Active integration will enable live training dashboard, hyperparameter sweeps, model versioning."),
    ]

    for i, (color, priority, title, desc) in enumerate(items):
        bx = Inches(0.4)
        by = Inches(1.3 + i * 1.16)
        bw = Inches(12.5)
        bh = Inches(1.05)
        add_rect(s, bx, by, bw, bh, BG_CARD)
        add_rect(s, bx, by, Inches(0.08), bh, color)

        # priority badge
        add_rect(s, bx + Inches(0.2), by + Inches(0.25),
                 Inches(0.85), Inches(0.4), color)
        add_text(s, priority,
                 bx + Inches(0.2), by + Inches(0.25),
                 Inches(0.85), Inches(0.4),
                 size=9, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

        add_text(s, title,
                 bx + Inches(1.2), by + Inches(0.1),
                 Inches(11.0), Inches(0.35),
                 size=14, bold=True, color=color)
        add_text(s, desc,
                 bx + Inches(1.2), by + Inches(0.5),
                 Inches(11.0), Inches(0.5),
                 size=11.5, color=GRAY)
    return s


def slide_16_conclusion(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_title(s, "Conclusion", "What FINQUANT-NEXUS v4 achieves")

    points = [
        (ACCENT_G, "Unified Intelligence Platform",
         "Single system combining RL, GNN, Federated Learning, and Sentiment NLP — no existing open-source tool integrates all four."),
        (ACCENT_B, "Real-World Data & Validation",
         "10 years of actual NIFTY 50 price data, live Google News feeds, and a true train/validate temporal split — not synthetic benchmarks."),
        (ACCENT_A, "Privacy by Design",
         "DP-SGD with (ε=8.0, δ=10⁻⁵) ensures collaborative learning without exposing any client's raw trade data — directly addresses SEBI constraints."),
        (ACCENT_R, "Risk-Aware, Not Just Return-Optimized",
         "Monte Carlo stress testing against 4 real crash scenarios (2008, COVID, Flash) — portfolio survives 82%+ under COVID conditions."),
        (RGBColor(0xA7,0x8B,0xFA), "Production-Ready Dashboard",
         "7 interactive tabs, 246 tests passing, Docker support, FastAPI backend — not a research prototype, a deployable platform."),
    ]

    for i, (color, title, desc) in enumerate(points):
        bx = Inches(0.4)
        by = Inches(1.3 + i * 1.16)
        add_rect(s, bx, by, Inches(12.5), Inches(1.05), BG_CARD)
        add_rect(s, bx, by, Inches(0.08), Inches(1.05), color)
        add_text(s, f"✓  {title}",
                 bx + Inches(0.2), by + Inches(0.1),
                 Inches(12.1), Inches(0.35),
                 size=14, bold=True, color=color)
        add_text(s, desc,
                 bx + Inches(0.2), by + Inches(0.5),
                 Inches(12.1), Inches(0.5),
                 size=12, color=GRAY)
    return s


def slide_17_thankyou(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, W, Inches(0.12), ACCENT_G)
    add_rect(s, 0, 0, Inches(0.08), H, ACCENT_G)

    add_text(s, "Thank You",
             Inches(1), Inches(1.8), Inches(11.3), Inches(1.1),
             size=52, bold=True, color=ACCENT_G, align=PP_ALIGN.CENTER)

    add_text(s, "Open for Questions & Live Demo",
             Inches(1), Inches(2.95), Inches(11.3), Inches(0.55),
             size=22, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    divider(s, Inches(3.7))

    add_text(s, "Praveen Rawal  ·  M.Tech DSML 2nd Year",
             Inches(1), Inches(3.9), Inches(11.3), Inches(0.45),
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, "FINQUANT-NEXUS v4  —  Self-Optimizing Federated Portfolio Intelligence Platform for NIFTY 50",
             Inches(1), Inches(4.45), Inches(11.3), Inches(0.45),
             size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # bottom quick stats
    stats = [
        ("1.87",  "Sharpe Ratio"),
        ("28.4%", "Annual Return"),
        ("82%",   "COVID Survival"),
        ("ε=8.0", "FL Privacy"),
        ("246",   "Tests Passing"),
        ("7",     "Dashboard Tabs"),
    ]
    for i, (val, label) in enumerate(stats):
        bx = Inches(0.5 + i * 2.05)
        add_rect(s, bx, Inches(5.6), Inches(1.9), Inches(1.4), BG_CARD)
        add_text(s, val,
                 bx, Inches(5.7), Inches(1.9), Inches(0.65),
                 size=22, bold=True, color=ACCENT_G, align=PP_ALIGN.CENTER)
        add_text(s, label,
                 bx, Inches(6.3), Inches(1.9), Inches(0.35),
                 size=10, color=GRAY, align=PP_ALIGN.CENTER)

    add_rect(s, 0, H - Inches(0.08), W, Inches(0.08), ACCENT_G)
    return s


def slide_18_industry(prs):
    s = blank_slide(prs)
    fill_bg(s)

    # Top accent bar (cyan — cyber theme)
    CYBER   = RGBColor(0x06, 0xB6, 0xD4)   # cyan-500
    CYBER2  = RGBColor(0x0E, 0x74, 0x90)   # deeper cyan
    add_rect(s, 0, 0, W, Inches(0.1), CYBER)

    # Section label
    add_text(s, "INDUSTRY EXPERIENCE",
             Inches(0.5), Inches(0.18), Inches(6), Inches(0.38),
             size=11, bold=True, color=CYBER)

    # Company name — large
    add_text(s, "Forensic Cybertech Pvt. Ltd.",
             Inches(0.5), Inches(0.6), Inches(9), Inches(0.85),
             size=34, bold=True, color=WHITE)

    # Role badge
    add_rect(s, Inches(0.5), Inches(1.45), Inches(2.2), Inches(0.44), CYBER)
    add_text(s, "AIML Intern",
             Inches(0.5), Inches(1.45), Inches(2.2), Inches(0.44),
             size=15, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

    add_text(s, "AI / ML  |  Cybersecurity Domain  |  Active",
             Inches(2.85), Inches(1.5), Inches(6), Inches(0.38),
             size=13, color=GRAY)

    divider(s, Inches(2.05))

    # ── Left column: What I do ───────────────────────────────────────────────
    add_rect(s, Inches(0.4), Inches(2.2), Inches(5.9), Inches(4.8), BG_CARD)
    add_rect(s, Inches(0.4), Inches(2.2), Inches(0.07), Inches(4.8), CYBER)

    add_text(s, "My Work at Forensic Cybertech",
             Inches(0.6), Inches(2.3), Inches(5.5), Inches(0.4),
             size=15, bold=True, color=CYBER)

    work_items = [
        ("AI Log Pipeline",
         "Design and build end-to-end pipelines that ingest, parse, and analyze cybersecurity logs at scale. Automates threat pattern detection that previously required manual analyst review."),
        ("AI Models for Cyber",
         "Train and deploy ML models for anomaly detection, malware classification, and intrusion detection on real network and system log data."),
        ("AI Agents",
         "Build autonomous agents that reason over security events, correlate indicators of compromise (IOCs), and recommend or execute response actions — enhancing SOC team efficiency."),
    ]

    for i, (title, desc) in enumerate(work_items):
        ty = Inches(2.85 + i * 1.35)
        add_rect(s, Inches(0.55), ty, Inches(0.38), Inches(0.38), CYBER)
        add_text(s, str(i + 1),
                 Inches(0.55), ty, Inches(0.38), Inches(0.38),
                 size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        add_text(s, title,
                 Inches(1.05), ty, Inches(5.1), Inches(0.35),
                 size=13, bold=True, color=WHITE)
        add_text(s, desc,
                 Inches(1.05), ty + Inches(0.38), Inches(5.1), Inches(0.85),
                 size=11, color=GRAY)

    # ── Right column: Impact cards ───────────────────────────────────────────
    add_rect(s, Inches(6.7), Inches(2.2), Inches(6.2), Inches(4.8), BG_CARD)
    add_rect(s, Inches(6.7), Inches(2.2), Inches(0.07), Inches(4.8), CYBER2)

    add_text(s, "Domain Focus & Impact",
             Inches(6.9), Inches(2.3), Inches(5.8), Inches(0.4),
             size=15, bold=True, color=CYBER)

    impact_cards = [
        (CYBER,
         "Efficiency Enhancement",
         "AI-driven automation reduces manual effort for security analysts — faster triage, fewer false positives, shorter incident response time."),
        (ACCENT_A,
         "Real-World Deployment",
         "Pipelines and models are built for production cybersecurity environments — handling live log streams and real threat intelligence feeds."),
        (ACCENT_G,
         "Bridging ML + Security",
         "Applies the same ML techniques from academic projects (RL agents, NLP, pipelines) directly to enterprise-grade cybersecurity problems."),
    ]

    for i, (color, title, desc) in enumerate(impact_cards):
        cy = Inches(2.85 + i * 1.35)
        add_rect(s, Inches(6.85), cy, Inches(5.9), Inches(1.2), RGBColor(0x0C, 0x1C, 0x2E))
        add_rect(s, Inches(6.85), cy, Inches(0.06), Inches(1.2), color)
        add_text(s, title,
                 Inches(7.05), cy + Inches(0.1), Inches(5.5), Inches(0.35),
                 size=13, bold=True, color=color)
        add_text(s, desc,
                 Inches(7.05), cy + Inches(0.48), Inches(5.5), Inches(0.65),
                 size=11, color=GRAY)

    # Bottom tagline
    add_rect(s, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3),
             RGBColor(0x06, 0x2A, 0x36))
    add_text(s,
             "Forensic Cybertech Pvt. Ltd.  —  Leveraging AI to strengthen cybersecurity operations and reduce analyst workload",
             Inches(0.6), Inches(7.1), Inches(12.1), Inches(0.3),
             size=11, italic=True, color=CYBER, align=PP_ALIGN.CENTER)

    add_rect(s, 0, H - Inches(0.08), W, Inches(0.08), CYBER)
    return s


# ─── MAIN ────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()

    print("Building slides...")
    slide_01_title(prs);     print("  01 Title")
    slide_02_agenda(prs);    print("  02 Agenda")
    slide_03_problem(prs);   print("  03 Problem Statement")
    slide_04_solution(prs);  print("  04 Solution Overview")
    slide_05_dataset(prs);   print("  05 Dataset")
    slide_06_to_12_tabs(prs);print("  06-12 Tab Slides (7 tabs)")
    slide_13_techstack(prs); print("  13 Tech Stack")
    slide_14_results(prs);   print("  14 Key Results")
    slide_15_future(prs);    print("  15 Future Work")
    slide_16_conclusion(prs);print("  16 Conclusion")
    slide_18_industry(prs);  print("  17 Industry Experience — Forensic Cybertech")
    slide_17_thankyou(prs);  print("  18 Thank You")

    out = "FINQUANT_NEXUS_Review2.pptx"
    prs.save(out)
    print(f"\nSaved: {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    build()
